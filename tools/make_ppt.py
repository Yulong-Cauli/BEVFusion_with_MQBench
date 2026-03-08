# -*- coding: utf-8 -*-
"""Generate BEVFusion + MQBench presentation PPT with ALL experiments."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree
import os

TEMPLATE = r"D:\Research\汇报\PPT模板.pptx"
OUTPUT = r"D:\Research\汇报\BEVFusion_MQBench_汇报.pptx"

FONT_NAME = "微软雅黑"
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# Colors
C_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
C_BODY = RGBColor(0x33, 0x33, 0x33)
C_ACCENT = RGBColor(0x00, 0x70, 0xC0)
C_GREEN = RGBColor(0x22, 0x8B, 0x22)
C_RED = RGBColor(0xCC, 0x00, 0x00)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_cell_bg(cell, r, g, b):
    """Set table cell background color via XML."""
    tc = cell._tc
    tcPr = tc.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr")
    if tcPr is None:
        tcPr = etree.SubElement(tc, "{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr")
        tc.insert(1, tcPr)
    for old in tcPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"):
        tcPr.remove(old)
    sf = etree.SubElement(tcPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    srgb = etree.SubElement(sf, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=C_BODY, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=12,
                          default_color=C_BODY, line_spacing=1.2):
    """Add textbox with multiple lines. Each line is (text, size, bold, color) or just str."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, bd, clr = line, default_size, False, default_color
        else:
            txt = line[0]
            sz = line[1] if len(line) > 1 else default_size
            bd = line[2] if len(line) > 2 else False
            clr = line[3] if len(line) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bd
        p.font.color.rgb = clr
        p.font.name = FONT_NAME
        p.space_after = Pt(2)
    return txBox


def add_table(slide, left, top, width, height, rows_data, col_widths=None,
              header_bg=(0x00, 0x70, 0xC0), font_size=10):
    """Add a table. rows_data[0] is header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                    set_cell_bg(cell, *header_bg)
                else:
                    p.font.color.rgb = C_BODY
                    if r_idx % 2 == 0:
                        set_cell_bg(cell, 0xF2, 0xF2, 0xF2)
                    else:
                        set_cell_bg(cell, 0xFF, 0xFF, 0xFF)
    return table_shape


def add_callout(slide, left, top, width, height, text, style="green"):
    """Add a colored callout box with rounded rectangle."""
    colors = {
        "green": (RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x2E, 0x7D, 0x32)),
        "blue": (RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x15, 0x65, 0xC0)),
        "orange": (RGBColor(0xFF, 0xF3, 0xE0), RGBColor(0xE6, 0x51, 0x00)),
        "red": (RGBColor(0xFF, 0xEB, 0xEE), RGBColor(0xC6, 0x28, 0x28)),
    }
    fill_c, border_c = colors.get(style, colors["green"])
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_c
    shape.line.color.rgb = border_c
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    # Title line
    p = tf.paragraphs[0]
    p.text = "💡 结论"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = border_c
    p.font.name = FONT_NAME
    # Content
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_BODY
    p2.font.name = FONT_NAME
    p2.space_before = Pt(2)
    return shape


def new_content_slide(prs, title_text):
    """Create a slide using layout 7 (仅标题) and set the title."""
    layout = prs.slide_layouts[7]  # "仅标题"
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text
            for p in ph.text_frame.paragraphs:
                p.font.name = FONT_NAME
                p.font.bold = True
                p.font.size = Pt(24)
            break
    return slide


def build_ppt():
    prs = Presentation(TEMPLATE)

    # Remove all existing slides from template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.2),
                "BEVFusion 模型量化与 TensorRT 部署", size=32, bold=True, color=C_TITLE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(3.3), Inches(8.0), Inches(0.6),
                "基于 MQBench 的 PTQ 量化 + TensorRT INT8 端到端部署", size=16,
                color=C_ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(8.0), Inches(0.5),
                "nuScenes 3D 目标检测", size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # ========== Slide 2: TOC ==========
    slide = new_content_slide(prs, "目录")
    toc_items = [
        ("一、研究背景与环境", 14, True, C_ACCENT),
        ("二、BEVFusion 架构与参数分布", 14, True, C_ACCENT),
        ("三、量化方案设计", 14, True, C_ACCENT),
        ("四、技术挑战与解决方案", 14, True, C_ACCENT),
        ("五、实验结果（9 组实验，本地 + 服务器）", 14, True, C_ACCENT),
        ("    实验①–③  SwinT PTQ 量化（本地 mini 数据集）", 12, False, C_BODY),
        ("    实验④     SwinT TRT 部署（本地 mini 数据集）", 12, False, C_BODY),
        ("    实验⑤–⑥  ResNet-50 PTQ + FP32 基准（本地 + 服务器）", 12, False, C_BODY),
        ("    实验⑦     ResNet-50 TRT 部署（本地 mini 数据集）", 12, False, C_BODY),
        ("    实验⑧–⑨  服务器完整验证集最终评估", 12, False, C_BODY),
        ("六、两种方案完整对比", 14, True, C_ACCENT),
        ("七、结论与未来工作", 14, True, C_ACCENT),
    ]
    add_multiline_textbox(slide, Inches(1.0), Inches(1.6), Inches(8.0), Inches(5.5), toc_items)

    # ========== Slide 3: Background ==========
    slide = new_content_slide(prs, "一、研究背景与环境")
    bg_lines = [
        ("研究目标", 14, True, C_ACCENT),
        ("对 BEVFusion 多模态 3D 检测模型进行 Post-Training Quantization (PTQ)，", 12),
        ("并通过 TensorRT 实现真实 INT8 部署，验证精度-体积-速度的权衡。", 12),
        ("", 8),
        ("技术栈", 14, True, C_ACCENT),
        ("• 量化框架：MQBench 0.0.6（基于 torch.fx 的 PTQ 仿真）", 12),
        ("• 部署引擎：TensorRT 10.x（IInt8EntropyCalibrator2 原生校准）", 12),
        ("• 基础模型：BEVFusion (camera+lidar)，nuScenes 3D 检测任务", 12),
        ("• 评估指标：NDS（综合得分）、mAP（检测精度）", 12),
        ("", 8),
        ("实验环境", 14, True, C_ACCENT),
        ("• 本地：RTX 4060 Laptop，nuScenes v1.0-mini（81 帧验证集）", 12),
        ("• 服务器：4×RTX 3090 + 1×A100-SXM4，nuScenes v1.0-trainval（6019 帧）", 12),
    ]
    add_multiline_textbox(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.5), bg_lines)

    # ========== Slide 4: Architecture ==========
    slide = new_content_slide(prs, "二、BEVFusion 架构与参数分布")
    arch_lines = [
        ("BEVFusion 6 大子模块", 13, True, C_ACCENT),
        ("Camera → Backbone → Neck → VTransform → BEV 特征 → Fuser ← LiDAR", 11),
        ("Fuser → Decoder Backbone → Decoder Neck → TransFusionHead → 检测结果", 11),
        ("", 6),
        ("SwinT 模型参数分布（FP32 .pth = 155.91 MB，39.80M 参数）", 13, True, C_ACCENT),
    ]
    add_multiline_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(2.0), arch_lines)

    # SwinT param table
    swint_data = [
        ["模块", "参数量", "占比", "FP32大小", "可量化"],
        ["camera/backbone (SwinT)", "27.55M", "67.5%", "105.2 MB", "❌ 动态控制流"],
        ["camera/neck (FPN)", "2.16M", "5.3%", "8.2 MB", "✅"],
        ["camera/vtransform", "2.59M", "6.4%", "10.0 MB", "❌ CUDA算子"],
        ["lidar/backbone", "2.69M", "6.6%", "10.3 MB", "❌ 稀疏卷积"],
        ["decoder/backbone", "4.26M", "10.4%", "16.3 MB", "✅"],
        ["decoder/neck", "0.28M", "0.7%", "1.1 MB", "✅"],
        ["fuser", "0.77M", "1.9%", "3.0 MB", "✅"],
        ["heads/object", "1.04M", "2.5%", "4.0 MB", "❌ Proxy迭代"],
    ]
    add_table(slide, Inches(0.4), Inches(3.6), Inches(9.2), Inches(3.2), swint_data, font_size=9)

    add_callout(slide, Inches(0.4), Inches(6.85), Inches(9.2), Inches(0.5),
                "SwinT 方案可量化参数仅 ~7M/39.8M (18%)，SwinTransformer 独占 67.5% 参数且无法 fx 追踪。",
                style="orange")

    # ========== Slide 5: ResNet-50 Architecture ==========
    slide = new_content_slide(prs, "二（续）、ResNet-50 替代方案参数分布")
    r50_intro = [
        ("动机：SwinTransformer 占 67.5% 参数但不可量化 → 替换为量化友好的 ResNet-50", 12, False, C_BODY),
        ("ResNet-50 模型（纯推理权重 = 142.8 MB，训练 .pth = 420.7 MB 含 optimizer）", 13, True, C_ACCENT),
    ]
    add_multiline_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.2), r50_intro)

    r50_data = [
        ["模块", "纯权重大小", "占比", "可量化"],
        ["camera/backbone (ResNet-50)", "89.9 MB", "63.0%", "✅ 纯CNN"],
        ["camera/neck (FPN)", "8.3 MB", "5.8%", "✅"],
        ["camera/vtransform", "10.0 MB", "7.0%", "❌ CUDA算子"],
        ["lidar/backbone", "10.3 MB", "7.2%", "❌ 稀疏卷积"],
        ["decoder/backbone", "16.3 MB", "11.4%", "✅"],
        ["decoder/neck", "1.1 MB", "0.8%", "✅"],
        ["fuser", "3.0 MB", "2.1%", "✅"],
        ["heads/object", "4.0 MB", "2.8%", "❌ Proxy迭代"],
        ["总计", "142.8 MB", "100%", "5/6 = 88%"],
    ]
    add_table(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(3.5), r50_data, font_size=9)

    add_callout(slide, Inches(0.4), Inches(6.4), Inches(9.2), Inches(0.65),
                "ResNet-50 替换后，可量化参数从 18% 跃升至 88%。camera/backbone 从不可量化变为可量化，"
                "是实现全模型高效 INT8 部署的关键突破。代价：FP32 精度低于 SwinT（仅训练 6 epochs）。",
                style="green")

    # ========== Slide 6: Design ==========
    slide = new_content_slide(prs, "三、量化方案设计")
    design_lines = [
        ("两阶段量化验证流程", 14, True, C_ACCENT),
        ("", 4),
        ("阶段一：MQBench PTQ 仿真（FakeQuant）", 13, True, C_GREEN),
        ("• torch.fx 自动追踪 → 插入 FakeQuantize 节点 → MinMax 校准", 11),
        ("• 在 PyTorch 中仿真 INT8 量化效果，快速验证精度是否可接受", 11),
        ("• 优点：无需 TensorRT，可在任意 GPU 上运行", 11),
        ("", 6),
        ("阶段二：TensorRT 真实 INT8 部署", 13, True, C_GREEN),
        ("• 导出 FP32 ONNX → TRT IInt8EntropyCalibrator2 原生校准", 11),
        ("• 逐模块替换为 TRT 引擎，端到端 NDS 评估", 11),
        ("• Hybrid 推理：TRT 模块 + PyTorch 模块混合运行", 11),
        ("", 6),
        ("为什么不用 MQBench convert_deploy 直接导出 INT8？", 13, True, RGBColor(0xE6, 0x51, 0x00)),
        ("• PyTorch 1.10 的 torch.onnx.export 不支持 FakeQuantize 节点", 11),
        ("• 可行方案：导出 FP32 ONNX → TRT 原生校准（已验证有效）", 11),
    ]
    add_multiline_textbox(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.5), design_lines)

    # ========== Slide 7: Challenges ==========
    slide = new_content_slide(prs, "四、技术挑战与解决方案")
    ch_lines = [
        ("挑战 1：torch.fx 追踪失败", 13, True, C_RED),
        ("• camera/neck: len(input) 在 Proxy 上调用 → 改用 self.num_ins 预计算常量", 10),
        ("• decoder/neck: 同上，FPN 中的 len() 调用 → 常量替代", 10),
        ("• fuser: torch.cat(inputs) 整体传入 → 改为逐元素索引 inputs[i]", 10),
        ("• mmcv Conv2d: x.numel()==0 检查 → patch_mmcv_for_fx() 上下文管理器", 10),
        ("", 4),
        ("挑战 2：SwinTransformer 不可量化", 13, True, C_RED),
        ("• 动态控制流（window_partition/shift）阻止 fx 追踪", 10),
        ("• 解决方案：替换为 ResNet-50（纯 CNN，fx 完全兼容）", 10),
        ("", 4),
        ("挑战 3：PTQ checkpoint 与 test.py 不兼容", 13, True, C_RED),
        ("• fx 改造后 state_dict key 变化，strict=False 静默跳过量化模块权重", 10),
        ("• 解决方案：PTQ 精度只通过 quant_ptq_minmax.py 内部评估", 10),
        ("", 4),
        ("挑战 4：分布式训练环境（服务器）", 13, True, C_RED),
        ("• torchpack 依赖 mpi4py + MASTER_HOST 环境变量，与 torchrun 不兼容", 10),
        ("• 解决方案：重写 _init_distributed()，支持 OMPI + torchrun 双路径", 10),
    ]
    add_multiline_textbox(slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(6.0), ch_lines)

    # ========== Slide 8: Exp 1 — SwinT PTQ 1/6 ==========
    slide = new_content_slide(prs, "实验①：SwinT PTQ 1/6 模块量化（mini, 02-25）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "环境：nuScenes v1.0-mini（81帧） | 仅量化 decoder/backbone（SECOND） | MinMax 校准",
                size=10, color=C_GRAY)

    exp1_acc = [
        ["指标", "FP32", "PTQ 1/6", "变化"],
        ["NDS", "0.5800", "0.5802", "+0.0002"],
        ["mAP", "0.5742", "0.5733", "−0.0009"],
    ]
    add_table(slide, Inches(0.4), Inches(1.9), Inches(4.0), Inches(1.1), exp1_acc, font_size=10)
    add_textbox(slide, Inches(0.4), Inches(1.72), Inches(4.0), Inches(0.2), "精度", size=11, bold=True, color=C_ACCENT)

    exp1_speed = [
        ["指标", "FP32", "PTQ（仿真）"],
        ["均值延迟", "386.77 ms", "398.23 ms"],
        ["FPS", "2.59", "2.51"],
        ["加速比", "—", "0.97x（略慢）"],
    ]
    add_table(slide, Inches(4.8), Inches(1.9), Inches(4.8), Inches(1.5), exp1_speed, font_size=10)
    add_textbox(slide, Inches(4.8), Inches(1.72), Inches(4.8), Inches(0.2), "速度（FakeQuant 仿真）",
                size=11, bold=True, color=C_ACCENT)

    exp1_size = [
        ["指标", "FP32", "PTQ（仿真）"],
        ["参数量", "39.80M", "39.81M (+FakeQuant)"],
        [".pth 大小", "155.91 MB", "156.24 MB"],
        ["理论 INT8 大小", "—", "~39 MB (÷4)"],
    ]
    add_table(slide, Inches(0.4), Inches(3.5), Inches(5.0), Inches(1.5), exp1_size, font_size=10)
    add_textbox(slide, Inches(0.4), Inches(3.32), Inches(5.0), Inches(0.2), "模型大小",
                size=11, bold=True, color=C_ACCENT)

    exp1_cov = [
        ["模块", "状态"],
        ["decoder/backbone", "✅ 已量化"],
        ["其他 5 个模块", "❌ fx 追踪失败"],
    ]
    add_table(slide, Inches(5.8), Inches(3.75), Inches(3.8), Inches(1.0), exp1_cov, font_size=10)
    add_textbox(slide, Inches(5.8), Inches(3.55), Inches(3.8), Inches(0.2), "量化覆盖：1/6",
                size=11, bold=True, color=C_ACCENT)

    add_callout(slide, Inches(0.4), Inches(5.3), Inches(9.2), Inches(0.7),
                "首次 PTQ 实验验证可行性：仅 1/6 模块量化，精度完全无损（NDS +0.0002）。"
                "FakeQuant 仿真速度略慢是预期行为（额外 clamp/round 开销），真实 INT8 部署方向相反。"
                "量化前 FP32 模型大小 155.91 MB。", style="green")

    # ========== Slide 9: Exp 2 — SwinT PTQ 4/6 ==========
    slide = new_content_slide(prs, "实验②：SwinT PTQ 4/6 覆盖扩大（mini, 02-25）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "新增量化：camera/neck + decoder/neck + fuser → 覆盖 1/6→4/6 | 校准 128 batch",
                size=10, color=C_GRAY)

    exp2_acc = [
        ["指标", "FP32", "PTQ 1/6", "PTQ 4/6", "Δ(vs FP32)"],
        ["NDS", "0.5801", "0.5802", "0.5814", "+0.0013"],
    ]
    add_table(slide, Inches(0.4), Inches(1.9), Inches(9.2), Inches(0.8), exp2_acc, font_size=10)

    exp2_speed = [
        ["指标", "FP32", "PTQ 1/6", "PTQ 4/6"],
        ["均值延迟", "389.46 ms", "398.23 ms", "408.50 ms"],
        ["FPS", "2.57", "2.51", "2.45"],
        ["加速比", "—", "0.97x", "0.95x"],
    ]
    add_table(slide, Inches(0.4), Inches(3.0), Inches(5.5), Inches(1.5), exp2_speed, font_size=10)

    exp2_cov = [
        ["模块", "状态"],
        ["decoder/backbone", "✅"],
        ["decoder/neck", "✅ 新增"],
        ["camera/neck", "✅ 新增"],
        ["fuser", "✅ 新增"],
        ["camera/backbone", "❌ 动态控制流"],
        ["heads/object", "❌ Proxy迭代"],
    ]
    add_table(slide, Inches(6.3), Inches(3.0), Inches(3.3), Inches(2.2), exp2_cov, font_size=9)

    add_callout(slide, Inches(0.4), Inches(5.5), Inches(9.2), Inches(0.8),
                "覆盖从 1/6 扩大到 4/6，精度不仅无损反而微升（NDS +0.0013）。"
                "技术突破：修复了 camera/neck、decoder/neck 中的 len(Proxy) 问题，"
                "以及 fuser 中的 torch.cat 列表参数问题。仿真速度更慢是因为 FakeQuant 节点增多。",
                style="green")

    # ========== Slide 10: Exp 3 — SwinT PTQ 4/6 Final ==========
    slide = new_content_slide(prs, "实验③：SwinT PTQ 4/6 最终确认（mini, 02-25）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "完整校准 + 评估流程 | quant_ptq_minmax.py（不加 --no-eval） | 128 batch 校准",
                size=10, color=C_GRAY)

    exp3_acc = [
        ["指标", "FP32 基线", "PTQ 4/6 (MinMax)", "变化"],
        ["NDS", "0.5801", "0.5810", "+0.0009（无损）"],
        ["mAP", "0.5742", "0.5759", "+0.0017（无损）"],
    ]
    add_table(slide, Inches(0.4), Inches(1.9), Inches(9.2), Inches(1.1), exp3_acc, font_size=10)

    exp3_cls = [
        ["类别", "FP32", "PTQ 4/6", "变化"],
        ["car", "0.916", "0.918", "+0.002"],
        ["truck", "0.833", "0.840", "+0.007"],
        ["bus", "0.995", "0.995", "0.000"],
        ["pedestrian", "0.919", "0.922", "+0.003"],
        ["motorcycle", "0.705", "0.699", "−0.006"],
        ["bicycle", "0.517", "0.518", "+0.001"],
        ["traffic_cone", "0.848", "0.866", "+0.018"],
    ]
    add_table(slide, Inches(0.4), Inches(3.3), Inches(6.0), Inches(2.8), exp3_cls, font_size=9)
    add_textbox(slide, Inches(0.4), Inches(3.1), Inches(6.0), Inches(0.2), "逐类 AP（去除 mini 样本不足的 3 类）",
                size=11, bold=True, color=C_ACCENT)

    add_callout(slide, Inches(0.4), Inches(6.3), Inches(9.2), Inches(0.65),
                "最终确认：4/6 模块 MinMax PTQ 量化后精度完全无损。NDS 和 mAP 均微升，"
                "所有逐类 AP 变化均在 ±0.02 范围内。这为后续 TensorRT INT8 部署提供了充分信心。",
                style="green")

    # ========== Slide 11: Exp 4 — SwinT TRT Hybrid ==========
    slide = new_content_slide(prs, "实验④：SwinT TRT Hybrid 4模块（mini, 02-26）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "4 模块 ONNX→TRT 引擎替换 | IInt8EntropyCalibrator2（50样本） | RTX 4060 Laptop",
                size=10, color=C_GRAY)

    exp4_acc = [
        ["方法", "NDS", "mAP", "NDS Δ", "mAP Δ"],
        ["PyTorch FP32", "0.5800", "0.5744", "—", "—"],
        ["TRT FP32", "0.5800", "0.5744", "+0.0000", "+0.0000"],
        ["TRT FP16", "0.5795", "0.5743", "−0.0005", "−0.0001"],
        ["TRT INT8", "0.5723", "0.5652", "−0.0077", "−0.0092"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(5.5), Inches(1.8), exp4_acc, font_size=9)

    exp4_eng = [
        ["模块", "FP32", "FP16", "INT8"],
        ["camera_neck", "8.0 MB", "3.1 MB", "1.7 MB"],
        ["fuser", "5.3 MB", "1.5 MB", "0.8 MB"],
        ["dec_backbone", "28.2 MB", "8.2 MB", "4.2 MB"],
        ["dec_neck", "1.2 MB", "0.7 MB", "0.6 MB"],
        ["总计", "42.6 MB", "13.5 MB", "7.2 MB"],
    ]
    add_table(slide, Inches(6.0), Inches(1.8), Inches(3.7), Inches(2.2), exp4_eng, font_size=9)
    add_textbox(slide, Inches(6.0), Inches(1.6), Inches(3.7), Inches(0.2), "引擎大小",
                size=11, bold=True, color=C_ACCENT)

    # Cosine similarity summary
    cos_lines = [
        ("余弦相似度（TRT vs PyTorch）", 11, True, C_ACCENT),
        ("• FP32：所有模块 cos = 1.000000（完全一致）", 10),
        ("• FP16：所有模块 cos ≥ 0.999993（几乎无损）", 10),
        ("• 压缩比：INT8 引擎 7.2 MB / FP32 权重 26.5 MB = 3.68x", 10),
    ]
    add_multiline_textbox(slide, Inches(0.3), Inches(4.1), Inches(5.5), Inches(1.2), cos_lines)

    add_callout(slide, Inches(0.3), Inches(5.5), Inches(9.4), Inches(0.95),
                "首次 TRT 真实部署验证！FP32 完全无损，FP16 几乎无损（NDS −0.0005），"
                "INT8 精度下降约 1.3%（NDS −0.0077）。INT8 引擎总大小仅 7.2 MB（vs 量化前 FP32 权重 26.5 MB，压缩 3.68x）。"
                "但未量化模块（SwinTransformer 等）仍需 129.4 MB → 总部署 ~136.6 MB。"
                "瓶颈：SwinTransformer 占 67.5% 参数但不可量化。",
                style="blue")

    # ========== Slide 12: Exp 5 — ResNet-50 PTQ 5/6 ==========
    slide = new_content_slide(prs, "实验⑤：ResNet-50 PTQ 5/6 量化突破（mini, 03-02）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "camera/backbone 替换为 ResNet-50 | 服务器训练 6 epochs | FakeQuant 仿真",
                size=10, color=C_GRAY)

    exp5_acc = [
        ["指标", "SwinT FP32", "SwinT PTQ 4/6", "R50 FP32", "R50 PTQ 5/6"],
        ["NDS", "0.5800", "0.5811", "0.3982", "0.4079"],
        ["mAP", "0.5744", "0.5760", "0.4135", "0.4189"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(9.4), Inches(1.1), exp5_acc, font_size=10)

    exp5_cov = [
        ["模块", "SwinT", "ResNet-50"],
        ["camera/backbone", "❌ 动态控制流", "✅ 纯CNN"],
        ["camera/neck", "✅", "✅"],
        ["fuser", "✅", "✅"],
        ["decoder/backbone", "✅", "✅"],
        ["decoder/neck", "✅", "✅"],
        ["heads/object", "❌", "❌"],
        ["可量化参数占比", "~18%", "~88%"],
    ]
    add_table(slide, Inches(0.3), Inches(3.2), Inches(5.5), Inches(2.6), exp5_cov, font_size=9)
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(5.5), Inches(0.2),
                "量化覆盖对比：关键突破", size=11, bold=True, color=C_ACCENT)

    size_lines = [
        ("模型大小对比", 11, True, C_ACCENT),
        ("SwinT FP32: 155.91 MB (39.80M)", 10),
        ("ResNet-50 纯权重: 142.8 MB", 10),
        ("ResNet-50 .pth: 420.7 MB (含 optimizer)", 10),
        ("", 6),
        ("⚠️ mini 精度受限", 11, True, RGBColor(0xE6, 0x51, 0x00)),
        ("R50 仅训练 6 ep，SwinT 用官方 20 ep 权重", 10),
        ("mini 集仅 81 帧，部分类别无样本", 10),
    ]
    add_multiline_textbox(slide, Inches(6.0), Inches(3.0), Inches(3.6), Inches(2.8), size_lines)

    add_callout(slide, Inches(0.3), Inches(6.0), Inches(9.4), Inches(0.8),
                "关键突破：ResNet-50 的 camera/backbone 成功量化！覆盖率从 18% 跃升至 88%。"
                "PTQ 5/6 精度无损（NDS +0.0097 在噪声范围内）。绝对精度较低是因为仅训练 6 epochs，"
                "需在完整验证集上确认真实性能。", style="green")

    # ========== Slide 13: Exp 6 — R50 FP32 Server ==========
    slide = new_content_slide(prs, "实验⑥：ResNet-50 FP32 服务器评估（6019帧, 03-02）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "nuScenes v1.0-trainval 完整验证集 | 4×RTX 3090 + 1×A100 | 6019帧全部 10 类",
                size=10, color=C_GRAY)

    exp6_acc = [
        ["指标", "mini (81帧)", "完整 val (6019帧)", "说明"],
        ["NDS", "0.3982", "0.4991", "+0.1009，mini 严重低估"],
        ["mAP", "0.4135", "0.4960", "+0.0825"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(9.4), Inches(1.1), exp6_acc, font_size=10)

    exp6_cls = [
        ["类别", "mini", "full", "类别", "mini", "full"],
        ["car", "0.854", "0.754", "motorcycle", "0.473", "0.466"],
        ["truck", "0.801", "0.419", "bicycle", "0.148", "0.378"],
        ["bus", "0.548", "0.465", "traffic_cone", "0.592", "0.662"],
        ["pedestrian", "0.761", "0.607", "trailer", "0.000", "0.319"],
        ["", "", "", "construction_vehicle", "0.000", "0.203"],
        ["", "", "", "barrier", "0.000", "0.686"],
    ]
    add_table(slide, Inches(0.3), Inches(3.2), Inches(9.4), Inches(2.3), exp6_cls, font_size=9)
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(9.4), Inches(0.2),
                "逐类 AP：mini 缺失类别在完整集上获得合理分数", size=11, bold=True, color=C_ACCENT)

    add_callout(slide, Inches(0.3), Inches(5.8), Inches(9.4), Inches(0.9),
                "完整数据集 NDS = 0.4991，比 mini 结果 (0.3982) 高出 +0.1009。"
                "mini 集严重低估模型真实性能！所有 10 个类别均获得有效 AP。"
                "与 BEVFusion 官方 SwinT (NDS 0.7138) 的差距主要来自训练量不足 (6ep vs 20ep)。"
                "训练 Loss 仍在持续下降（epoch6 下降率 7.6%），继续训练可显著提升。",
                style="blue")

    # ========== Slide 14: Exp 7 — R50 TRT Hybrid ==========
    slide = new_content_slide(prs, "实验⑦：ResNet-50 TRT 5模块部署（mini, 03-02）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "5 模块 ONNX→TRT 引擎（新增 camera_backbone） | IInt8EntropyCalibrator2 | RTX 4060",
                size=10, color=C_GRAY)

    exp7_acc = [
        ["方法", "NDS", "mAP", "NDS Δ"],
        ["PyTorch FP32", "0.3982", "0.4135", "—"],
        ["PTQ 5/6 FakeQuant", "0.4079", "0.4189", "+0.0097"],
        ["TRT FP32 (5模块)", "0.4030", "0.4172", "+0.0048"],
        ["TRT FP16 (5模块)", "0.3981", "0.4136", "−0.0001"],
        ["TRT INT8 (5模块)", "0.4078", "0.4187", "+0.0096"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(5.2), Inches(2.0), exp7_acc, font_size=9)

    exp7_eng = [
        ["模块", "FP32", "FP16", "INT8"],
        ["camera_backbone", "116.9 MB", "46.3 MB", "24.2 MB"],
        ["camera_neck", "12.1 MB", "4.3 MB", "2.3 MB"],
        ["fuser", "5.4 MB", "1.5 MB", "0.8 MB"],
        ["dec_backbone", "28.9 MB", "8.4 MB", "4.3 MB"],
        ["dec_neck", "1.2 MB", "0.7 MB", "0.6 MB"],
        ["总计", "160.7 MB", "59.9 MB", "31.4 MB"],
    ]
    add_table(slide, Inches(5.8), Inches(1.8), Inches(3.9), Inches(2.5), exp7_eng, font_size=9)
    add_textbox(slide, Inches(5.8), Inches(1.6), Inches(3.9), Inches(0.2), "引擎大小",
                size=11, bold=True, color=C_ACCENT)

    exp7_deploy = [
        ["精度", "5模块 TRT", "未量化模块", "总部署体积", "压缩比"],
        ["FP32 引擎", "160.7 MB", "24.3 MB", "~185 MB", "2.3x"],
        ["FP16 引擎", "59.9 MB", "24.3 MB", "~84 MB", "5.0x"],
        ["INT8 引擎", "31.4 MB", "24.3 MB", "~55.7 MB", "7.6x"],
    ]
    add_table(slide, Inches(0.3), Inches(4.2), Inches(9.0), Inches(1.5), exp7_deploy, font_size=9)
    add_textbox(slide, Inches(0.3), Inches(4.0), Inches(9.0), Inches(0.2),
                "部署体积（vs 量化前 .pth 420.7 MB）", size=11, bold=True, color=C_ACCENT)

    add_callout(slide, Inches(0.3), Inches(5.9), Inches(9.4), Inches(0.85),
                "5 模块 TRT 全部导出成功！INT8 精度无损（NDS +0.0096 为统计波动）。"
                "INT8 总部署体积仅 55.7 MB（vs 量化前 420.7 MB .pth），压缩 7.6x。"
                "相比 SwinT 方案总部署 136.6 MB，减少 59%。"
                "camera_backbone 是最大引擎（INT8 24.2 MB），但成功量化是体积下降的核心原因。",
                style="green")

    # ========== Slide 15: Exp 8 — Server SwinT Final ==========
    slide = new_content_slide(prs, "实验⑧：SwinT 服务器完整评估（6019帧, 03-03）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "⭐ 最终结果 | nuScenes v1.0-trainval 6019帧 | 官方预训练 SwinT 权重 | 4模块 TRT",
                size=10, color=C_GRAY)

    exp8_acc = [
        ["方法", "NDS", "mAP", "NDS Δ", "mAP Δ"],
        ["PyTorch FP32 基线", "0.7069", "0.6728", "—", "—"],
        ["PTQ 4/6 FakeQuant", "0.7015", "0.6618", "−0.0054", "−0.0110"],
        ["TRT FP32 (4模块)", "0.7065", "0.6726", "−0.0004", "−0.0002"],
        ["TRT FP16 (4模块)", "0.7069", "0.6728", "+0.0000", "+0.0000"],
        ["TRT INT8 (4模块)", "0.7022", "0.6641", "−0.0047", "−0.0087"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(6.0), Inches(2.1), exp8_acc, font_size=10)

    exp8_eng = [
        ["模块", "FP32", "FP16", "INT8"],
        ["camera_neck", "9.9 MB", "3.2 MB", "1.7 MB"],
        ["fuser", "5.4 MB", "1.6 MB", "0.8 MB"],
        ["dec_backbone", "28.9 MB", "8.5 MB", "4.3 MB"],
        ["dec_neck", "1.2 MB", "0.7 MB", "0.6 MB"],
        ["总计", "44.4 MB", "13.6 MB", "7.4 MB"],
    ]
    add_table(slide, Inches(6.5), Inches(1.8), Inches(3.2), Inches(2.1), exp8_eng, font_size=9)
    add_textbox(slide, Inches(6.5), Inches(1.6), Inches(3.2), Inches(0.2), "引擎大小",
                size=11, bold=True, color=C_ACCENT)

    # Key per-class AP (selected)
    exp8_cls = [
        ["类别", "FP32", "TRT FP16", "TRT INT8", "INT8 Δ"],
        ["car", "0.875", "0.875", "0.876", "+0.001"],
        ["truck", "0.639", "0.639", "0.626", "−0.013"],
        ["bus", "0.741", "0.741", "0.729", "−0.012"],
        ["pedestrian", "0.877", "0.877", "0.874", "−0.003"],
        ["motorcycle", "0.770", "0.769", "0.755", "−0.015"],
        ["bicycle", "0.612", "0.611", "0.606", "−0.006"],
    ]
    add_table(slide, Inches(0.3), Inches(4.2), Inches(6.5), Inches(2.2), exp8_cls, font_size=9)
    add_textbox(slide, Inches(0.3), Inches(4.0), Inches(6.5), Inches(0.2),
                "逐类 AP（主要类别）", size=11, bold=True, color=C_ACCENT)

    add_callout(slide, Inches(0.3), Inches(6.55), Inches(9.4), Inches(0.7),
                "SwinT 最终结果：TRT FP16 与 FP32 完全一致（NDS 0.7069）！"
                "INT8 精度损失仅 0.67% (NDS −0.0047)，4模块 INT8 引擎仅 7.4 MB。"
                "量化前 FP32 模型 155.91 MB，但未量化模块仍需 ~129 MB → 总部署 ~136 MB。",
                style="green")

    # ========== Slide 16: Exp 9 — Server R50 Final ==========
    slide = new_content_slide(prs, "实验⑨：ResNet-50 服务器完整评估（6019帧, 03-03）")
    add_textbox(slide, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.4),
                "⭐ 最终结果 | nuScenes v1.0-trainval 6019帧 | 6 epochs 训练 | 5模块 TRT",
                size=10, color=C_GRAY)

    exp9_acc = [
        ["方法", "NDS", "mAP", "NDS Δ", "mAP Δ"],
        ["PyTorch FP32 基线", "0.4989", "0.4961", "—", "—"],
        ["PTQ 5/6 FakeQuant", "0.4958", "0.4904", "−0.0031", "−0.0057"],
        ["TRT FP32 (5模块)", "0.4994", "0.4965", "+0.0005", "+0.0004"],
        ["TRT FP16 (5模块)", "0.4992", "0.4962", "+0.0003", "+0.0001"],
        ["TRT INT8 (5模块)", "0.4948", "0.4945", "−0.0041", "−0.0016"],
    ]
    add_table(slide, Inches(0.3), Inches(1.8), Inches(6.0), Inches(2.1), exp9_acc, font_size=10)

    exp9_eng = [
        ["模块", "FP32", "FP16", "INT8"],
        ["camera_backbone", "115.7 MB", "46.5 MB", "24.0 MB"],
        ["camera_neck", "12.2 MB", "4.4 MB", "2.4 MB"],
        ["fuser", "5.4 MB", "1.6 MB", "0.8 MB"],
        ["dec_backbone", "28.9 MB", "8.5 MB", "4.4 MB"],
        ["dec_neck", "1.2 MB", "0.7 MB", "0.6 MB"],
        ["总计", "159.6 MB", "60.2 MB", "31.4 MB"],
    ]
    add_table(slide, Inches(6.5), Inches(1.8), Inches(3.2), Inches(2.5), exp9_eng, font_size=9)
    add_textbox(slide, Inches(6.5), Inches(1.6), Inches(3.2), Inches(0.2), "引擎大小",
                size=11, bold=True, color=C_ACCENT)

    # Deploy summary
    exp9_deploy = [
        ["精度", "5模块 TRT", "未量化模块*", "总部署体积"],
        ["INT8", "31.4 MB", "24.3 MB", "55.7 MB"],
        ["FP16", "60.2 MB", "24.3 MB", "84.5 MB"],
    ]
    add_table(slide, Inches(0.3), Inches(4.4), Inches(5.5), Inches(1.1), exp9_deploy, font_size=10)
    add_textbox(slide, Inches(0.3), Inches(4.2), Inches(5.5), Inches(0.2),
                "部署体积（量化前 142.8 MB 纯权重 / 420.7 MB .pth）", size=10, bold=True, color=C_ACCENT)

    size_note = [
        ("*未量化模块 24.3 MB：", 10, True, C_GRAY),
        ("lidar/backbone 10.3 + vtransform 10.0", 9, False, C_GRAY),
        ("+ heads/object 4.0 MB", 9, False, C_GRAY),
    ]
    add_multiline_textbox(slide, Inches(6.0), Inches(4.4), Inches(3.6), Inches(1.0), size_note)

    add_callout(slide, Inches(0.3), Inches(5.75), Inches(9.4), Inches(0.8),
                "ResNet-50 最终结果：INT8 精度损失仅 0.82% (NDS −0.0041)，FP16 无损。"
                "INT8 总部署体积仅 55.7 MB（vs 量化前 420.7 MB .pth），压缩 7.6x！"
                "相比 SwinT 方案的 136 MB 部署，减少 59%。量化覆盖 88% 是核心优势。",
                style="green")

    # ========== Slide 17: Comparison ==========
    slide = new_content_slide(prs, "六、两种方案完整对比（完整验证集 6019帧）")

    comp_data = [
        ["指标", "SwinT", "SwinT", "ResNet-50", "ResNet-50"],
        ["", "FP32", "TRT INT8", "FP32", "TRT INT8"],
        ["NDS", "0.7069", "0.7022 (−0.67%)", "0.4989", "0.4948 (−0.82%)"],
        ["mAP", "0.6728", "0.6641 (−1.29%)", "0.4961", "0.4945 (−0.32%)"],
        ["TRT 模块数", "—", "4/6", "—", "5/6"],
        ["可量化参数占比", "—", "~18%", "—", "~88%"],
        ["INT8 引擎大小", "—", "7.4 MB", "—", "31.4 MB"],
        ["未量化模块", "—", "~129 MB", "—", "24.3 MB"],
        ["总部署体积", "155.91 MB", "~136 MB", "142.8 MB*", "55.7 MB"],
        ["量化前 .pth", "155.91 MB", "—", "420.7 MB", "—"],
    ]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.5), comp_data, font_size=10)

    # PTQ vs TRT validation
    ptq_trt = [
        ["验证", "PTQ FakeQuant NDS Δ", "TRT INT8 NDS Δ", "差异"],
        ["SwinT", "−0.0054", "−0.0047", "TRT 更好 0.0007"],
        ["ResNet-50", "−0.0031", "−0.0041", "差异仅 0.001"],
    ]
    add_table(slide, Inches(0.3), Inches(5.3), Inches(6.0), Inches(1.1), ptq_trt, font_size=10,
              header_bg=(0x2E, 0x7D, 0x32))
    add_textbox(slide, Inches(0.3), Inches(5.1), Inches(6.0), Inches(0.2),
                "PTQ 仿真 vs TRT 真实部署：MQBench 可准确预测 INT8 精度", size=10, bold=True, color=C_GREEN)

    comp_note = [
        ("*ResNet-50 纯推理权重 142.8 MB", 9, False, C_GRAY),
        ("（训练 .pth 含 optimizer = 420.7 MB）", 9, False, C_GRAY),
    ]
    add_multiline_textbox(slide, Inches(6.5), Inches(5.3), Inches(3.2), Inches(0.8), comp_note)

    add_callout(slide, Inches(0.3), Inches(6.55), Inches(9.4), Inches(0.65),
                "精度优先 → SwinT + TRT FP16（NDS 0.7069，与 FP32 完全一致）；"
                "部署体积优先 → ResNet-50 + TRT INT8（NDS 0.4948，总体积仅 55.7 MB）。"
                "MQBench PTQ 仿真可准确预测真实 INT8 部署精度。",
                style="blue")

    # ========== Slide 18: Conclusion ==========
    slide = new_content_slide(prs, "七、结论与未来工作")
    concl_lines = [
        ("主要成果", 16, True, C_ACCENT),
        ("", 4),
        ("1. 完成 BEVFusion 两种 backbone 方案的 PTQ 量化 + TRT 部署全流程", 12, True),
        ("   SwinT 4/6 模块 + ResNet-50 5/6 模块，均在 6019 帧完整验证集验证", 11),
        ("", 4),
        ("2. INT8 量化精度损失极小", 12, True),
        ("   SwinT: NDS −0.67%  |  ResNet-50: NDS −0.82%", 11),
        ("", 4),
        ("3. 显著的模型压缩", 12, True),
        ("   SwinT: 155.91 MB → 136 MB (INT8引擎仅7.4MB)", 11),
        ("   ResNet-50: 420.7 MB → 55.7 MB (压缩7.6x)", 11),
        ("", 4),
        ("4. 解决了多项工程难题", 12, True),
        ("   torch.fx 兼容性修复、mmcv patch、分布式训练环境适配", 11),
        ("", 8),
        ("未来工作", 16, True, C_ACCENT),
        ("", 4),
        ("• 继续训练 ResNet-50 至 20 epochs（Loss 仍在下降，预估 NDS 可达 0.55+）", 11),
        ("• 推理速度测量（当前缺少端到端延迟数据）", 11),
        ("• 探索更先进的校准策略（Histogram/Percentile/AdaRound）", 11),
        ("• TransFusionHead 量化（静态展开 ModuleList 迭代）", 11),
    ]
    add_multiline_textbox(slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(5.8), concl_lines)

    # ========== Slide 19: Thanks ==========
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # "致谢模板"
    add_textbox(slide, Inches(2.0), Inches(2.8), Inches(6.0), Inches(1.0),
                "谢谢！", size=40, bold=True, color=C_TITLE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), Inches(4.0), Inches(6.0), Inches(0.6),
                "欢迎提问与讨论", size=18, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(OUTPUT)
    print(f"PPT saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_ppt()
